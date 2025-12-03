import os
from typing import List, Dict, Optional

import docx2txt
from PyPDF2 import PdfReader
from pptx import Presentation

from config import DATA_DIR


class DocumentLoader:
    def __init__(
        self,
        data_dir: str = DATA_DIR,
    ):
        self.data_dir = data_dir
        self.supported_formats = [".pdf", ".pptx", ".docx", ".txt"]

    def load_pdf(self, file_path: str) -> List[Dict]:
        """加载PDF文件，按页返回内容

        TODO: 实现PDF文件加载
        要求：
        1. 使用PdfReader读取PDF文件
        2. 遍历每一页，提取文本内容
        3. 格式化为"--- 第 X 页 ---\n文本内容\n"
        4. 返回pdf内容列表，每个元素包含 {"text": "..."}
        """
        pdf_content = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    formatted_text = f"--- 第 {page_num + 1} 页 ---\n{text}\n"
                    pdf_content.append({"text": formatted_text})
            return pdf_content
        except FileNotFoundError:
            print(f"错误：文件 {file_path} 未找到")
        except PdfReader.PdfReadError:
            print(f"错误：无法读取PDF文件 {file_path}")
        except Exception as e:
            print(f"读取PDF时发生错误：{e}")
    

    def load_pptx(self, file_path: str) -> List[Dict]:
        """加载PPT文件，按幻灯片返回内容

        要求：
        1. 使用Presentation读取PPT文件
        2. 遍历每一页，提取所有文本内容
        3. 格式化为"--- 幻灯片 X ---\n文本内容\n"
        4. 返回幻灯片内容列表，每个元素包含 {"text": "..."}
        """
        pptx_content = []
        try:
            presentation = Presentation(file_path)
            for slide_num in range(len(presentation.slides)):
                slide = presentation.slides[slide_num]

                # 收集当前幻灯片的所有文本
                slide_texts = []

                # 遍历所有形状，提取文本
                for shape in slide.shapes:
                    # 只处理有text属性的形状
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                # 如果幻灯片有文本内容，合并所有文本
                if slide_texts:
                    combined_text = "\n".join(slide_texts)
                    formatted_text = f"--- 幻灯片 {slide_num + 1} ---\n{combined_text}\n"
                    pptx_content.append({"text": formatted_text})
                else:
                    # 如果没有文本内容，也创建一个空条目
                    formatted_text = f"--- 幻灯片 {slide_num + 1} ---\n[无文本内容]\n"
                    pptx_content.append({"text": formatted_text})

            return pptx_content
        except FileNotFoundError:
            print(f"错误：文件 {file_path} 未找到")
        except (IOError, OSError) as e:
            print(f"错误：无法读取PPT文件 {file_path} - {e}")
        except Exception as e:
            print(f"读取PPT时发生错误：{e}")
        return pptx_content

    def load_docx(self, file_path: str) -> str:
        """加载DOCX文件
        TODO: 实现DOCX文件加载
        要求：
        1. 使用docx2txt读取DOCX文件
        2. 返回文本内容
        """
        docx_content = []
        try:
            text = docx2txt.process(file_path)
            formatted_text = f"--- DOCX文件 ---\n{text}\n"
            docx_content.append({"text": formatted_text})
        except FileNotFoundError:
            print(f"错误：文件 {file_path} 未找到")
        except docx2txt.Docx2txtError:
            print(f"错误：无法读取DOCX文件 {file_path}")
        except Exception as e:
            print(f"读取DOCX时发生错误：{e}")
        return docx_content

    def load_txt(self, file_path: str) -> str:
        """加载TXT文件
        TODO: 实现TXT文件加载
        要求：
        1. 使用open读取TXT文件（注意使用encoding="utf-8"）
        2. 返回文本内容
        """
        txt_content = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
                formatted_text = f"--- TXT文件 ---\n{text}\n"
                txt_content.append({"text": formatted_text})
        except FileNotFoundError:
            print(f"错误：文件 {file_path} 未找到")
        except UnicodeDecodeError:
            print(f"错误：无法读取TXT文件 {file_path}")
        except Exception as e:
            print(f"读取TXT时发生错误：{e}")
        return txt_content

    def load_document(self, file_path: str) -> List[Dict[str, str]]:
        """加载单个文档，PDF和PPT按页/幻灯片分割，返回文档块列表"""
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        documents = []

        if ext == ".pdf":
            pages = self.load_pdf(file_path)
            for page_idx, page_data in enumerate(pages, 1):
                documents.append(
                    {
                        "content": page_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": page_idx,
                    }
                )
        elif ext == ".pptx":
            slides = self.load_pptx(file_path)
            for slide_idx, slide_data in enumerate(slides, 1):
                documents.append(
                    {
                        "content": slide_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": slide_idx,
                    }
                )
        elif ext == ".docx":
            content = self.load_docx(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        elif ext == ".txt":
            content = self.load_txt(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        else:
            print(f"不支持的文件格式: {ext}")

        return documents

    def load_all_documents(self) -> List[Dict[str, str]]:
        """加载数据目录下的所有文档"""
        if not os.path.exists(self.data_dir):
            print(f"数据目录不存在: {self.data_dir}")
            return None

        documents = []

        for root, dirs, files in os.walk(self.data_dir):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in self.supported_formats:
                    file_path = os.path.join(root, file)
                    print(f"正在加载: {file_path}")
                    doc_chunks = self.load_document(file_path)
                    if doc_chunks:
                        documents.extend(doc_chunks)

        return documents
